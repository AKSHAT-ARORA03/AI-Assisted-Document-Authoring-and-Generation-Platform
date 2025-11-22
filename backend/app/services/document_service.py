from docx import Document
from docx.shared import Inches, Pt
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.style import WD_STYLE_TYPE
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from typing import List, Dict, Any
import io
import os

class DocumentService:
    @staticmethod
    def generate_word_document(project_data: Dict[str, Any]) -> io.BytesIO:
        """Generate a Word document from project data"""
        doc = Document()
        
        # Set up styles
        styles = doc.styles
        
        # Title style
        title_style = styles['Title']
        title_style.font.size = Pt(24)
        
        # Heading style
        heading_style = styles['Heading 1']
        heading_style.font.size = Pt(16)
        
        # Add title
        title = doc.add_heading(project_data['title'], 0)
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER
        
        # Add topic/description
        if project_data.get('description'):
            doc.add_paragraph(project_data['description'])
        
        doc.add_paragraph()  # Add space
        
        # Add sections
        sections = project_data.get('sections', [])
        sections.sort(key=lambda x: x.get('order', 0))  # Sort by order
        
        for section in sections:
            # Add section heading
            doc.add_heading(section['title'], level=1)
            
            # Add section content
            if section.get('content'):
                content = section['content']
                # Split content into paragraphs if it contains line breaks
                paragraphs = content.split('\n\n')
                for paragraph in paragraphs:
                    if paragraph.strip():
                        doc.add_paragraph(paragraph.strip())
            else:
                doc.add_paragraph("[Content not yet generated]")
            
            doc.add_paragraph()  # Add space between sections
        
        # Save to bytes
        buffer = io.BytesIO()
        doc.save(buffer)
        buffer.seek(0)
        return buffer
    
    @staticmethod
    def generate_powerpoint_presentation(project_data: Dict[str, Any]) -> io.BytesIO:
        """Generate a PowerPoint presentation from project data"""
        prs = Presentation()
        
        # Slide 1 - Title slide
        title_slide_layout = prs.slide_layouts[0]  # Title slide layout
        title_slide = prs.slides.add_slide(title_slide_layout)
        
        title = title_slide.shapes.title
        subtitle = title_slide.placeholders[1]
        
        title.text = project_data['title']
        if project_data.get('description'):
            subtitle.text = project_data['description']
        else:
            subtitle.text = f"A presentation about {project_data['topic']}"
        
        # Add content slides
        slides = project_data.get('slides', [])
        slides.sort(key=lambda x: x.get('order', 0))  # Sort by order
        
        for slide_data in slides:
            # Use title and content layout
            slide_layout = prs.slide_layouts[1]  # Title and content layout
            slide = prs.slides.add_slide(slide_layout)
            
            # Set title
            title = slide.shapes.title
            title.text = slide_data['title']
            
            # Add content
            content = slide.placeholders[1]
            tf = content.text_frame
            
            if slide_data.get('content'):
                bullet_points = slide_data['content']
                if isinstance(bullet_points, list):
                    for i, point in enumerate(bullet_points):
                        if i == 0:
                            tf.text = point
                        else:
                            p = tf.add_paragraph()
                            p.text = point
                            p.level = 0
                else:
                    tf.text = bullet_points
            else:
                tf.text = "[Content not yet generated]"
        
        # Save to bytes
        buffer = io.BytesIO()
        prs.save(buffer)
        buffer.seek(0)
        return buffer
    
    @staticmethod
    def get_file_extension(document_type: str) -> str:
        """Get file extension based on document type"""
        return "docx" if document_type == "docx" else "pptx"
    
    @staticmethod
    def get_mime_type(document_type: str) -> str:
        """Get MIME type based on document type"""
        if document_type == "docx":
            return "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
        else:
            return "application/vnd.openxmlformats-officedocument.presentationml.presentation"